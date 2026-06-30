from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_project_ppt(output_path):
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[6]
    cover_slide = prs.slides.add_slide(slide_layout)
    
    title_box = cover_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_text = title_box.text_frame
    title_text.word_wrap = True
    p = title_text.add_paragraph()
    p.text = "华南师范大学计算机专业"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = title_text.add_paragraph()
    p.text = "课程管理系统"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    subtitle_box = cover_slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    subtitle_text = subtitle_box.text_frame
    subtitle_text.word_wrap = True
    p = subtitle_text.add_paragraph()
    p.text = "基于检索增强生成（RAG）架构和多智能体协同机制"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    info_box = cover_slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(1))
    info_text = info_box.text_frame
    p = info_text.add_paragraph()
    p.text = "项目演示"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    
    bg_shape = cover_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(26, 35, 126)
    bg_shape.line.fill.background()
    
    slide_layout = prs.slide_layouts[1]
    toc_slide = prs.slides.add_slide(slide_layout)
    title = toc_slide.shapes.title
    title.text = "目录"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = toc_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    items = [
        "项目概述",
        "功能特性",
        "技术架构",
        "项目结构",
        "核心模块",
        "多Agent系统",
        "RAG技术实现",
        "总结与展望"
    ]
    
    for i, item in enumerate(items, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.level = 0
    
    slide_layout = prs.slide_layouts[1]
    overview_slide = prs.slides.add_slide(slide_layout)
    title = overview_slide.shapes.title
    title.text = "项目概述"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = overview_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "项目定位"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    
    p = tf.add_paragraph()
    p.text = "为华南师范大学计算机专业学生提供课程信息管理、电子资源访问、智能知识检索及多智能体协同对话等功能的智能课程管理系统。"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "核心技术"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "检索增强生成（RAG）架构 + 多智能体协同机制"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    slide_layout = prs.slide_layouts[1]
    features_slide = prs.slides.add_slide(slide_layout)
    title = features_slide.shapes.title
    title.text = "功能特性"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = features_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    features = [
        ("📚 课程管理", "课程列表展示、课程详情、资料下载、课程搜索"),
        ("📖 电子书库", "课程关联推荐、专业分类检索、关键词搜索"),
        ("🔍 智能知识检索", "RAG架构问答、流式输出、引用来源追溯"),
        ("🤖 多智能体协同对话", "理解Agent、检索Agent、推理Agent、生成Agent、验证Agent、推荐Agent"),
        ("⚙️ LLM模式设置", "知识检索模式、对话模式")
    ]
    
    for icon, desc in features:
        p = tf.add_paragraph()
        p.text = icon
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 35, 126)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.level = 1
    
    slide_layout = prs.slide_layouts[1]
    tech_slide = prs.slides.add_slide(slide_layout)
    title = tech_slide.shapes.title
    title.text = "技术架构"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = tech_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    tech_stack = [
        ("前端框架", "Vue 3 + Pinia + Tailwind CSS"),
        ("后端框架", "FastAPI + Uvicorn"),
        ("Agent框架", "LangGraph"),
        ("向量数据库", "Chroma"),
        ("业务数据库", "TiDB Cloud"),
        ("LLM接口", "SiliconFlow")
    ]
    
    for category, tech in tech_stack:
        p = tf.add_paragraph()
        p.text = f"• {category}: {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
    
    slide_layout = prs.slide_layouts[1]
    arch_slide = prs.slides.add_slide(slide_layout)
    title = arch_slide.shapes.title
    title.text = "架构层次"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = arch_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    layers = [
        ("表现层 (Presentation)", "Vue3组件、前端路由、UI界面"),
        ("应用层 (Application)", "认证服务、问答服务、知识库服务、学习服务"),
        ("领域层 (Domain)", "用户领域、知识领域、对话领域、学习领域"),
        ("基础设施层 (Infrastructure)", "数据库、向量库、LLM客户端、文件存储")
    ]
    
    y_pos = Inches(1.5)
    for layer_name, description in layers:
        shape = arch_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(8), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(33, 150, 243)
        shape.line.fill.background()
        
        text_box = arch_slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.2), Inches(7.6), Inches(0.6))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = layer_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        y_pos += Inches(1.2)
    
    slide_layout = prs.slide_layouts[1]
    structure_slide = prs.slides.add_slide(slide_layout)
    title = structure_slide.shapes.title
    title.text = "项目结构"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = structure_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "frontend/ - 前端项目"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    
    p = tf.add_paragraph()
    p.text = "├── src/components/ - Vue组件"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "├── src/stores/ - Pinia状态管理"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "backend/ - 后端项目"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "├── agents/ - Agent模块"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "├── application/ - 应用层服务"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "├── infrastructure/ - 基础设施层"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "├── interfaces/ - API接口层"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "└── domain/ - 领域层"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    
    slide_layout = prs.slide_layouts[1]
    modules_slide = prs.slides.add_slide(slide_layout)
    title = modules_slide.shapes.title
    title.text = "核心模块"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = modules_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    modules = [
        ("认证服务", "JWT令牌生成、用户注册登录、权限验证"),
        ("问答服务", "知识检索、流式响应、会话管理"),
        ("知识库服务", "向量检索、RAG实现、知识管理"),
        ("学习服务", "学习记录、进度追踪、个性化推荐")
    ]
    
    for name, desc in modules:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 35, 126)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.level = 1
    
    slide_layout = prs.slide_layouts[1]
    agent_slide = prs.slides.add_slide(slide_layout)
    title = agent_slide.shapes.title
    title.text = "多Agent系统架构"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = agent_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    agents = [
        ("🧠 UnderstandingAgent", "理解用户意图和实体"),
        ("🔍 RetrievalAgent", "从知识库检索相关知识"),
        ("🤔 ReasoningAgent", "基于知识进行逻辑推理"),
        ("✍️ GenerationAgent", "生成自然语言回答"),
        ("✅ ValidationAgent", "验证答案质量"),
        ("💡 RecommendAgent", "推荐相关知识")
    ]
    
    for icon, desc in agents:
        p = tf.add_paragraph()
        p.text = f"• {icon}: {desc}"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(60, 60, 60)
    
    slide_layout = prs.slide_layouts[1]
    rag_slide = prs.slides.add_slide(slide_layout)
    title = rag_slide.shapes.title
    title.text = "RAG技术实现"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    
    content = rag_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    rag_steps = [
        "1. 用户提问 → 语义理解",
        "2. 查询向量化 → Embedding模型",
        "3. 向量检索 → Chroma向量数据库",
        "4. 结果重排序 → Rerank模型",
        "5. 上下文构建 → 检索结果拼接",
        "6. 答案生成 → LLM推理",
        "7. 流式输出 → Server-Sent Events"
    ]
    
    for step in rag_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
    
    slide_layout = prs.slide_layouts[6]
    summary_slide = prs.slides.add_slide(slide_layout)
    
    title_box = summary_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_text = title_box.text_frame
    title_text.word_wrap = True
    p = title_text.add_paragraph()
    p.text = "总结与展望"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    p.alignment = PP_ALIGN.CENTER
    
    content_box = summary_slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    content_text = content_box.text_frame
    content_text.word_wrap = True
    
    p = content_text.add_paragraph()
    p.text = "已实现功能：课程管理、电子书库、智能问答、多Agent协同对话"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = PP_ALIGN.CENTER
    
    p = content_text.add_paragraph()
    p.text = "未来方向：个性化学习路径、知识图谱构建、移动端适配"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = PP_ALIGN.CENTER
    
    bg_shape = summary_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    bg_shape.line.fill.background()
    
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")

if __name__ == "__main__":
    output_dir = os.path.join(os.path.dirname(__file__), "..")
    output_path = os.path.join(output_dir, "项目介绍.pptx")
    create_project_ppt(output_path)